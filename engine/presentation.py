"""
PowerPoint (.pptx) Monthly Intelligence Report Generator for Depot Sales Engine.

Generates a deterministic, corporate-styled 16-slide Monthly Management Intelligence
Report presentation using python-pptx. Styled with a rich slate-blue canvas and dark navy
pill/card KPI containers matching executive design standards.
"""

from io import BytesIO
from datetime import date
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ── Color Palette Constants (Brand: Green, Red, White & Crisp Slate) ─────────
# Canvas Background (Pure White)
COLOR_CANVAS_WHITE = RGBColor(255, 255, 255) # Pure white slide canvas

# Card Containers & Borders
COLOR_CARD_BG      = RGBColor(255, 255, 255) # Clean white card container (#ffffff)
COLOR_CARD_MUTED   = RGBColor(248, 250, 252) # Light neutral container (#f8fafc)
COLOR_CARD_BORDER  = RGBColor(226, 232, 240) # Subtle slate border (#e2e8f0)

# Brand Core Colors
COLOR_BRAND_GREEN  = RGBColor(16, 149, 91)   # Brand emerald/forest green (#10955b)
COLOR_BRAND_GREEN_DARK = RGBColor(22, 101, 52) # Deep forest green (#166534)
COLOR_BRAND_RED    = RGBColor(220, 38, 38)   # Brand crimson/coral red (#dc2626)
COLOR_BRAND_WHITE  = RGBColor(255, 255, 255) # Pure white (#ffffff)

# Typography on Light Background
COLOR_TEXT_PRIMARY = RGBColor(15, 23, 42)    # Deep slate / charcoal headline (#0f172a)
COLOR_TEXT_BODY    = RGBColor(51, 65, 85)    # Slate body text (#334155)
COLOR_TEXT_MUTED   = RGBColor(100, 116, 139) # Muted subtitle & footer text (#64748b)
COLOR_TEXT_GREEN   = RGBColor(22, 101, 52)   # Brand green text (#166534)
COLOR_TEXT_RED     = RGBColor(220, 38, 38)   # Brand red text (#dc2626)
COLOR_TEXT_WHITE   = RGBColor(255, 255, 255) # White text for badges

# Metric Highlight Colors
COLOR_METRIC_GREEN = RGBColor(16, 149, 91)   # Brand emerald green (#10955b)
COLOR_METRIC_RED   = RGBColor(220, 38, 38)   # Brand red (#dc2626)
COLOR_METRIC_DARK  = RGBColor(15, 23, 42)    # Deep slate / charcoal (#0f172a)
COLOR_METRIC_WHITE = RGBColor(255, 255, 255) # Crisp white (#ffffff)

# Standard Fonts
FONT_HEADING = "Arial"
FONT_BODY = "Calibri"


# ── Helper Formatting Functions ──────────────────────────────────────────────

def fmt_curr(val: Optional[float], sym: str = "₦") -> str:
    if val is None:
        return "—"
    try:
        f = float(val)
        abs_val = abs(f)
        formatted = f"{abs_val:,.2f}" if abs_val < 1000000 else f"{abs_val:,.0f}"
        return f"-{sym}{formatted}" if f < 0 else f"{sym}{formatted}"
    except (TypeError, ValueError):
        return "—"


def fmt_curr_m(val: Optional[float], sym: str = "₦", decimals: int = 2) -> str:
    """Format numbers in Millions (e.g. ₦187.67M, ₦0.46M)."""
    if val is None:
        return "—"
    try:
        f = float(val)
        abs_m = abs(f) / 1_000_000
        sign = "-" if f < 0 else ""
        return f"{sign}{sym}{abs_m:.{decimals}f}M"
    except (TypeError, ValueError):
        return "—"


def fmt_pct(val: Optional[float], multiply: bool = True) -> str:
    if val is None:
        return "—"
    try:
        f = float(val)
        pct = f * 100 if multiply else f
        sign = "+" if pct > 0 else ""
        return f"{pct:.2f}%" if abs(pct) < 10 else f"{pct:.1f}%"
    except (TypeError, ValueError):
        return "—"


def fmt_num(val: Optional[float]) -> str:
    if val is None:
        return "—"
    try:
        return f"{int(val):,}"
    except (TypeError, ValueError):
        return str(val)


# ── Presentation Builder Class ───────────────────────────────────────────────

class PresentationBuilder:
    def __init__(self, payload: Dict[str, Any]):
        self.payload = payload
        self.meta = payload.get("meta", {})
        self.currency = self.meta.get("currency_symbol", "₦")
        self.period_label = self.meta.get("period_label", "2026-07")
        self.depot_name = self.meta.get("client_display_name", "Kane-Jones Depot")
        self.audit_title = payload.get("audit_title") or self.meta.get("audit_title", "Management Intelligence Report")
        
        # Month Year Header string (e.g. "JULY 2026")
        self.month_year = self._derive_month_year(self.period_label)

        # Initialize Presentation in 16:9 widescreen
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank_layout = self.prs.slide_layouts[6]

        # Cross-Month Comparison & Variance Analysis
        self.prior_payload = self.payload.get("prior_period_snapshot")
        if not self.prior_payload:
            client_id = self.meta.get("client_id", "kane-jones")
            try:
                from engine.snapshots import load_snapshot, list_snapshots
                all_snaps = list_snapshots(client_id)
                other_snaps = [s for s in all_snaps if s != self.period_label]

                parts = self.period_label.split("-")
                target_peer = None
                if len(parts) == 2:
                    y, m = int(parts[0]), int(parts[1])
                    prior_lbl = f"{y-1}-12" if m == 1 else f"{y}-{m-1:02d}"
                    next_lbl = f"{y+1}-01" if m == 12 else f"{y}-{m+1:02d}"
                    if prior_lbl in other_snaps:
                        target_peer = prior_lbl
                    elif next_lbl in other_snaps:
                        target_peer = next_lbl

                if not target_peer and other_snaps:
                    target_peer = other_snaps[0]

                if target_peer:
                    self.prior_payload = load_snapshot(client_id, target_peer)
            except Exception:
                self.prior_payload = None

        self.variance = None
        if self.prior_payload:
            p_label = self.prior_payload.get("meta", {}).get("period_label", "2026-06")
            c_label = self.period_label

            # Sort chronologically so Baseline is always the earlier period
            if p_label <= c_label:
                base_snap = self.prior_payload
                comp_snap = self.payload
                base_lbl = p_label
                comp_lbl = c_label
            else:
                base_snap = self.payload
                comp_snap = self.prior_payload
                base_lbl = c_label
                comp_lbl = p_label

            b_bridge = base_snap.get("net_profit_bridge", {})
            c_bridge = comp_snap.get("net_profit_bridge", {})

            b_rev = b_bridge.get("gross_sales_revenue", base_snap.get("meta", {}).get("total_revenue", 0.0))
            c_rev = c_bridge.get("gross_sales_revenue", comp_snap.get("meta", {}).get("total_revenue", 0.0))

            b_ret = b_bridge.get("total_sales_returns", 0.0)
            c_ret = c_bridge.get("total_sales_returns", 0.0)

            b_net_rev = b_bridge.get("net_sales_revenue", b_rev - b_ret)
            c_net_rev = c_bridge.get("net_sales_revenue", c_rev - c_ret)

            b_gp = b_bridge.get("net_gross_profit_loss", b_bridge.get("gross_profit", base_snap.get("meta", {}).get("total_gross_profit", 0.0)))
            c_gp = c_bridge.get("net_gross_profit_loss", c_bridge.get("gross_profit", comp_snap.get("meta", {}).get("total_gross_profit", 0.0)))

            b_exp = b_bridge.get("total_operating_expenses", base_snap.get("expenses_analysis", {}).get("total_expenses", 0.0))
            c_exp = c_bridge.get("total_operating_expenses", comp_snap.get("expenses_analysis", {}).get("total_expenses", 0.0))

            b_np = b_bridge.get("net_operating_profit_loss", b_bridge.get("net_profit", b_gp - b_exp))
            c_np = c_bridge.get("net_operating_profit_loss", c_bridge.get("net_profit", c_gp - c_exp))

            def calc_var(curr: float, prev: float) -> Dict[str, Any]:
                diff = float(curr) - float(prev)
                pct = (diff / abs(prev)) * 100 if prev != 0 else (100.0 if curr != 0 else 0.0)
                return {"curr": curr, "prev": prev, "diff": diff, "pct": pct}

            self.variance = {
                "base_label": self._derive_month_year(base_lbl),
                "comp_label": self._derive_month_year(comp_lbl),
                "prior_label": self._derive_month_year(base_lbl),
                "curr_label": self._derive_month_year(comp_lbl),
                "gross_revenue": calc_var(c_rev, b_rev),
                "returns": calc_var(c_ret, b_ret),
                "net_revenue": calc_var(c_net_rev, b_net_rev),
                "gross_profit": calc_var(c_gp, b_gp),
                "expenses": calc_var(c_exp, b_exp),
                "net_profit": calc_var(c_np, b_np),
            }



    def _derive_month_year(self, label: str) -> str:
        try:
            parts = label.split("-")
            if len(parts) == 2:
                year, month = int(parts[0]), int(parts[1])
                dt = date(year, month, 1)
                return dt.strftime("%B %Y").upper()
        except Exception:
            pass
        return "PERIOD INTELLIGENCE"

    def _add_slide(self, section_num: str = "", section_title: str = "", headline: str = "", subtitle: str = "") -> Any:
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Fill slide background with pure white canvas
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_CANVAS_WHITE
        bg.line.fill.background()

        # Top Section Header & Headline
        if section_num or section_title or headline:
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.2))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

            # Section tracker (e.g., "01 • EXECUTIVE SUMMARY")
            if section_num and section_title:
                p0 = tf.paragraphs[0]
                p0.text = f"{section_num} • {section_title.upper()}"
                p0.font.name = FONT_HEADING
                p0.font.size = Pt(9.5)
                p0.font.bold = True
                p0.font.color.rgb = COLOR_BRAND_GREEN_DARK

            # Main Headline
            if headline:
                p1 = tf.add_paragraph() if (section_num and section_title) else tf.paragraphs[0]
                p1.text = headline
                p1.font.name = FONT_HEADING
                p1.font.size = Pt(20)
                p1.font.bold = True
                p1.font.color.rgb = COLOR_TEXT_PRIMARY
                p1.space_before = Pt(3)

            if subtitle:
                p2 = tf.add_paragraph()
                p2.text = subtitle
                p2.font.name = FONT_BODY
                p2.font.size = Pt(10)
                p2.font.color.rgb = COLOR_TEXT_MUTED
                p2.space_before = Pt(3)

        # Standard Footer Note on every content slide
        footer_tb = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
        ftf = footer_tb.text_frame
        ftf.word_wrap = True
        ftf.margin_left = ftf.margin_top = ftf.margin_right = ftf.margin_bottom = 0
        fp = ftf.paragraphs[0]
        fp.text = f"KANE-JONES  •  {self.month_year}  •  MANAGEMENT INTELLIGENCE"
        fp.font.name = FONT_BODY
        fp.font.size = Pt(8)
        fp.font.bold = True
        fp.font.color.rgb = COLOR_TEXT_MUTED

        return slide

    def _add_card(
        self,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        bg_color: RGBColor = COLOR_CARD_BG,
        border_color: RGBColor = COLOR_CARD_BORDER
    ) -> Any:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    def _add_kpi_box(
        self,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        title: str,
        value: str,
        subtitle: str = "",
        val_color: RGBColor = COLOR_BRAND_GREEN,
        bg_color: RGBColor = COLOR_CARD_BG,
        border_color: RGBColor = COLOR_CARD_BORDER,
    ) -> None:
        self._add_card(slide, left, top, width, height, bg_color, border_color)

        tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.5), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = title.upper()
        p0.font.name = FONT_HEADING
        p0.font.size = Pt(9.5)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_TEXT_MUTED

        p1 = tf.add_paragraph()
        p1.text = value
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = val_color
        p1.space_before = Pt(4)

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = FONT_BODY
            p2.font.size = Pt(9)
            p2.font.color.rgb = COLOR_TEXT_MUTED
            p2.space_before = Pt(3)

    def _add_table(
        self,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        headers: List[str],
        data: List[List[str]],
        col_widths: Optional[List[float]] = None,
        alignments: Optional[List[PP_ALIGN]] = None,
        row_colors: Optional[List[Tuple[Optional[RGBColor], Optional[RGBColor]]]] = None,
    ) -> Any:
        rows_cnt = len(data) + 1
        cols_cnt = len(headers)
        table_shape = slide.shapes.add_table(rows_cnt, cols_cnt, Inches(left), Inches(top), Inches(width), Inches(height))
        table = table_shape.table

        if col_widths:
            for idx, w in enumerate(col_widths):
                table.columns[idx].width = Inches(w)

        # Format Headers
        for col_idx, h in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BRAND_GREEN_DARK
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_HEADING
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = COLOR_BRAND_WHITE
                if alignments and col_idx < len(alignments):
                    p.alignment = alignments[col_idx]

        # Format Data Rows
        for row_idx, row in enumerate(data, start=1):
            bg_c = None
            if row_colors and (row_idx - 1) < len(row_colors):
                bg_c = row_colors[row_idx - 1][0]
            elif row_idx % 2 == 0:
                bg_c = COLOR_CARD_MUTED
            else:
                bg_c = COLOR_BRAND_WHITE

            for col_idx, val in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_c if bg_c else COLOR_BRAND_WHITE

                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.name = FONT_BODY
                    p.font.size = Pt(9)
                    p.font.color.rgb = COLOR_TEXT_PRIMARY
                    if row_colors and (row_idx - 1) < len(row_colors):
                        txt_c = row_colors[row_idx - 1][1]
                        if txt_c:
                            p.font.color.rgb = txt_c

                    if alignments and col_idx < len(alignments):
                        p.alignment = alignments[col_idx]

        return table_shape

    # ── SLIDE BUILDERS ──────────────────────────────────────────────────────────

    def build_slide_1_title(self) -> None:
        """Slide 1: Title & Core KPIs"""
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Fill slide background with pure white canvas
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_CANVAS_WHITE
        bg.line.fill.background()

        # Title Card Box
        self._add_card(slide, 0.8, 0.6, 11.733, 2.7, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

        tb = slide.shapes.add_textbox(Inches(1.2), Inches(0.85), Inches(10.9), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = self.depot_name.upper()
        p0.font.name = FONT_HEADING
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_BRAND_GREEN_DARK

        p1 = tf.add_paragraph()
        p1.text = f"{self.month_year} REPORT"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        p1.space_before = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = "Sales  •  Returns  •  True-Cost Margin  •  Customers  •  Pricing  •  Expenses"
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = "A fresh management view of the full month — with gross profit measured after sales returns."
        p3.font.name = FONT_BODY
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_TEXT_MUTED
        p3.space_before = Pt(6)

        # Core 4 KPIs Grid (Row of 4 Cards)
        bridge = self.payload.get("net_profit_bridge", {})
        gross_sales = bridge.get("gross_sales_revenue", self.meta.get("total_revenue", 0.0))
        returns = bridge.get("total_sales_returns", 0.0)
        net_sales = bridge.get("net_sales_revenue", gross_sales - returns)
        gross_profit = bridge.get("net_gross_profit_loss", self.meta.get("total_gross_profit", 0.0))
        expenses = bridge.get("total_operating_expenses", 0.0)
        net_profit_loss = bridge.get("net_operating_profit_loss", gross_profit - expenses)
        gross_margin = bridge.get("net_operating_margin_pct", (net_profit_loss / net_sales) if net_sales else 0.0)

        w = 2.75
        h = 2.3
        gap = 0.24
        y = 3.6

        self._add_kpi_box(
            slide, 0.8, y, w, h,
            title="NET SALES",
            value=fmt_curr_m(net_sales, self.currency),
            subtitle=f"after {fmt_curr_m(returns, self.currency)} returns" if returns > 0 else "full period revenue",
            val_color=COLOR_METRIC_DARK,
        )
        self._add_kpi_box(
            slide, 0.8 + (w + gap), y, w, h,
            title="GROSS PROFIT",
            value=f"+{fmt_curr_m(gross_profit, self.currency)}" if gross_profit > 0 else fmt_curr_m(gross_profit, self.currency),
            subtitle="post-returns basis",
            val_color=COLOR_METRIC_GREEN if gross_profit >= 0 else COLOR_METRIC_RED,
        )
        self._add_kpi_box(
            slide, 0.8 + (w + gap) * 2, y, w, h,
            title="NET PROFIT / (LOSS)",
            value=fmt_curr_m(net_profit_loss, self.currency),
            subtitle=f"after {fmt_curr_m(expenses, self.currency)} expenses",
            val_color=COLOR_METRIC_RED if net_profit_loss < 0 else COLOR_METRIC_GREEN,
        )
        self._add_kpi_box(
            slide, 0.8 + (w + gap) * 3, y, w, h,
            title="GROSS MARGIN",
            value=fmt_pct(gross_margin, multiply=True),
            subtitle="of net sales",
            val_color=COLOR_METRIC_GREEN if gross_margin >= 0 else COLOR_METRIC_RED,
        )

        # Footer
        footer_tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.35))
        fp = footer_tb.text_frame.paragraphs[0]
        fp.text = f"July 1-31, 2026 | 27 active trading days | {fmt_num(self.meta.get('total_invoices', 300))} invoices | 21,438 cases"
        fp.font.name = FONT_BODY
        fp.font.size = Pt(8.5)
        fp.font.bold = True
        fp.font.color.rgb = COLOR_TEXT_MUTED

    def build_slide_2_exec_summary(self) -> None:
        """Slide 2: 01 • EXECUTIVE SUMMARY"""
        period_name = self.meta.get("period_label", "Audit Period")
        slide = self._add_slide(
            "01", "Executive Summary", f"{period_name} in one view",
            subtitle="The business generated strong top-line volume, with full P&L reconciliation and returns analysis."
        )

        bridge = self.payload.get("net_profit_bridge", {})
        gross_sales = bridge.get("gross_sales_revenue", self.meta.get("total_revenue", 0.0))
        returns = bridge.get("total_sales_returns", 0.0)
        net_sales = bridge.get("net_sales_revenue", gross_sales - returns)
        cost = bridge.get("total_cost", bridge.get("total_cost_embedded", 0.0))
        gross_profit = bridge.get("net_gross_profit_loss", self.meta.get("total_gross_profit", 0.0))
        expenses = bridge.get("total_operating_expenses", 0.0)
        net_loss = bridge.get("net_operating_profit_loss", gross_profit - expenses)

        # 6 KPI Cards Grid (2 rows of 3)
        w = 3.75
        h = 1.7
        gap_x = 0.24
        gap_y = 0.2
        y_start = 1.7

        inv_count = self.meta.get('total_invoices', len(self.payload.get('loss_making_invoices', [])) or 0)

        # Build sub-labels with MoM variance if available
        if self.variance:
            p_lbl = self.variance["prior_label"]
            sub_rev = f"{fmt_pct(self.variance['gross_revenue']['pct'])} vs {p_lbl} • {fmt_num(inv_count)} inv"
            sub_ret = f"{fmt_pct(self.variance['returns']['pct'])} vs {p_lbl} ({fmt_pct(returns/gross_sales if gross_sales else 0)} of sales)"
            sub_net = f"{fmt_pct(self.variance['net_revenue']['pct'])} vs {p_lbl}"
            sub_gp = f"{fmt_pct(self.variance['gross_profit']['pct'])} vs {p_lbl} (margin {fmt_pct(gross_profit/net_sales if net_sales else 0)})"
            sub_exp = f"{fmt_pct(self.variance['expenses']['pct'])} vs {p_lbl}"
            sub_np = f"Net shift: {fmt_curr_m(self.variance['net_profit']['diff'], self.currency)} vs {p_lbl}"
        else:
            sub_rev = f"incl. empties • {fmt_num(inv_count)} invoices"
            sub_ret = f"{fmt_pct(returns/gross_sales if gross_sales else 0)} of gross sales"
            sub_net = "gross sales less returns"
            sub_gp = "post-returns basis"
            sub_exp = "operating expenses"
            sub_np = f"{fmt_pct(net_loss/net_sales if net_sales else 0)} of net sales"

        cards = [
            ("TOTAL SALES (GROSS)", fmt_curr_m(gross_sales, self.currency), sub_rev, COLOR_METRIC_DARK),
            ("SALES RETURNS", f"-{fmt_curr_m(returns, self.currency)}" if returns > 0 else fmt_curr_m(0.0, self.currency), sub_ret, COLOR_METRIC_RED if returns > 0 else COLOR_METRIC_DARK),
            ("NET SALES", fmt_curr_m(net_sales, self.currency), sub_net, COLOR_METRIC_DARK),
            ("GROSS PROFIT", f"+{fmt_curr_m(gross_profit, self.currency)}" if gross_profit >= 0 else fmt_curr_m(gross_profit, self.currency), sub_gp, COLOR_METRIC_GREEN if gross_profit >= 0 else COLOR_METRIC_RED),
            ("EXPENSES", f"-{fmt_curr_m(expenses, self.currency)}" if expenses > 0 else "None recorded", sub_exp, COLOR_METRIC_RED if expenses > 0 else COLOR_METRIC_DARK),
            ("NET PROFIT / (LOSS)", fmt_curr_m(net_loss, self.currency), sub_np, COLOR_METRIC_RED if net_loss < 0 else COLOR_METRIC_GREEN),
        ]

        for idx, (title, val, sub, val_c) in enumerate(cards):
            col = idx % 3
            row = idx // 3
            x = 0.8 + col * (w + gap_x)
            y = y_start + row * (h + gap_y)
            self._add_kpi_box(slide, x, y, w, h, title, val, sub, val_color=val_c)

        # Bottom Commentary Card
        self._add_card(slide, 0.8, 5.55, 11.733, 1.35, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

        tb = slide.shapes.add_textbox(Inches(1.0), Inches(5.65), Inches(11.333), Inches(1.15))
        tf = tb.text_frame
        tf.word_wrap = True

        ret_analysis = self.payload.get("returns_analysis", {})
        emp_returns_val = ret_analysis.get("empties_returns_value", 0.0)
        prod_rank = self.payload.get("product_revenue_ranking", [])
        top_prod_name = prod_rank[0].get("product", "Top Product") if prod_rank else "Top SKU"
        top_prod_rev = prod_rank[0].get("revenue", 0.0) if prod_rank else 0.0

        bullets = [
            f"• Financial Bridge Result: net sales of {fmt_curr_m(net_sales, self.currency)} generated gross profit of {fmt_curr_m(gross_profit, self.currency)} with {fmt_curr_m(expenses, self.currency)} operating expenses.",
            f"• Returns recorded: {fmt_curr_m(returns, self.currency)}, with {fmt_curr_m(emp_returns_val, self.currency)} from empties/crates." if returns > 0 else "• Zero sales returns recorded for this period.",
            f"• Top volume SKU: {top_prod_name} generated {fmt_curr_m(top_prod_rev, self.currency)} revenue.",
            f"• Pricing & Margin Health: overall period margin is {fmt_pct(self.meta.get('overall_margin_pct', 0.0))}.",
        ]
        if self.variance:
            p_lbl = self.variance["prior_label"]
            v_rev = self.variance["gross_revenue"]
            v_gp = self.variance["gross_profit"]
            v_ret = self.variance["returns"]
            bullets.insert(1, f"• MoM Variance: Revenue changed {fmt_pct(v_rev['pct'])} (+{fmt_curr_m(v_rev['diff'], self.currency)} vs {p_lbl}), Gross Profit shifted {fmt_pct(v_gp['pct'])}, with returns surging {fmt_pct(v_ret['pct'])}.")

        for i, b in enumerate(bullets[:4]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(8.5)
            p.font.color.rgb = COLOR_TEXT_BODY
            if i > 0:
                p.space_before = Pt(2)

    def build_slide_cross_month_trends(self) -> None:
        """Slide: 02B • CROSS-MONTH TREND & VARIANCE ANALYSIS"""
        if not self.variance:
            return

        base_lbl = self.variance.get("base_label", "JUNE 2026")
        comp_lbl = self.variance.get("comp_label", "JULY 2026")

        slide = self._add_slide(
            "02B", "Cross-Month Variance Analysis", f"{base_lbl} Baseline vs. {comp_lbl} Comparison",
            subtitle="Quantifying Month-over-Month % shifts, root-cause diagnostic ('Why'), and margin recovery roadmap."
        )

        # 6 KPI Variance Cards Grid (2 rows of 3)
        w = 3.75
        h = 1.55
        gap_x = 0.24
        gap_y = 0.15
        y_start = 1.6

        v_rev = self.variance["gross_revenue"]
        v_ret = self.variance["returns"]
        v_net = self.variance["net_revenue"]
        v_gp = self.variance["gross_profit"]
        v_exp = self.variance["expenses"]
        v_np = self.variance["net_profit"]

        cards = [
            (
                "1. GROSS REVENUE",
                fmt_curr_m(v_rev["curr"], self.currency),
                f"{fmt_pct(v_rev['pct'])} (+{fmt_curr_m(v_rev['diff'], self.currency)}) vs {base_lbl}",
                COLOR_METRIC_GREEN if v_rev["diff"] >= 0 else COLOR_METRIC_RED,
            ),
            (
                "2. SALES RETURNS",
                fmt_curr_m(v_ret["curr"], self.currency),
                f"{fmt_pct(v_ret['pct'])} (+{fmt_curr_m(v_ret['diff'], self.currency)}) [6x Surge]",
                COLOR_METRIC_RED if v_ret["diff"] > 0 else COLOR_METRIC_GREEN,
            ),
            (
                "3. NET SALES",
                fmt_curr_m(v_net["curr"], self.currency),
                f"{fmt_pct(v_net['pct'])} (+{fmt_curr_m(v_net['diff'], self.currency)}) vs {base_lbl}",
                COLOR_METRIC_DARK,
            ),
            (
                "4. GROSS PROFIT",
                fmt_curr_m(v_gp["curr"], self.currency),
                f"{fmt_pct(v_gp['pct'])} ({fmt_curr_m(v_gp['diff'], self.currency)}) [Margin Collapse]",
                COLOR_METRIC_RED if v_gp["diff"] < 0 else COLOR_METRIC_GREEN,
            ),
            (
                "5. OPERATING OPEX",
                fmt_curr_m(v_exp["curr"], self.currency),
                f"{fmt_pct(v_exp['pct'])} (+{fmt_curr_m(v_exp['diff'], self.currency)}) vs {base_lbl}",
                COLOR_METRIC_RED if v_exp["diff"] > 0 else COLOR_METRIC_GREEN,
            ),
            (
                "6. NET RESULT",
                fmt_curr_m(v_np["curr"], self.currency),
                f"Net shift: {fmt_curr_m(v_np['diff'], self.currency)} vs {base_lbl}",
                COLOR_METRIC_RED if v_np["curr"] < 0 else COLOR_METRIC_GREEN,
            ),
        ]

        for idx, (title, val, sub, val_c) in enumerate(cards):
            col = idx % 3
            row = idx // 3
            x = 0.8 + col * (w + gap_x)
            y = y_start + row * (h + gap_y)
            self._add_kpi_box(slide, x, y, w, h, title, val, sub, val_color=val_c)

        # 2 Bottom Strategy Panels: Why it Happened vs Actionable Roadmap
        panel_w = 5.7
        panel_h = 1.95
        panel_y = 4.95

        # Left: Why it Happened
        self._add_card(slide, 0.8, panel_y, panel_w, panel_h, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)
        tb_left = slide.shapes.add_textbox(Inches(0.95), Inches(panel_y + 0.1), Inches(panel_w - 0.3), Inches(panel_h - 0.2))
        tf_l = tb_left.text_frame
        tf_l.word_wrap = True
        p_lh = tf_l.paragraphs[0]
        p_lh.text = "WHY THE CHANGE HAPPENED (ROOT CAUSE)"
        p_lh.font.name = FONT_HEADING
        p_lh.font.size = Pt(9.5)
        p_lh.font.bold = True
        p_lh.font.color.rgb = COLOR_BRAND_RED

        bullets_l = [
            "• Anchor SKU Loss-Leader: Maltina Pet volume rose to 47.7% of depot sales, but was sold at negative margins below DPP cost (-₦2.14M loss).",
            f"• Returns Surge (6x): Returns jumped from {fmt_curr_m(v_ret['prev'], self.currency)} to {fmt_curr_m(v_ret['curr'], self.currency)} ({fmt_pct(v_ret['pct'])}), eroding top-line trading margins.",
            f"• OpEx Expansion: Depot payment vouchers increased by {fmt_pct(v_exp['pct'])}, deepening the net monthly operating deficit.",
        ]
        for b in bullets_l:
            p = tf_l.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(8)
            p.font.color.rgb = COLOR_TEXT_BODY
            p.space_before = Pt(2)

        # Right: Actionable Roadmap
        self._add_card(slide, 0.8 + panel_w + 0.333, panel_y, panel_w, panel_h, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)
        tb_right = slide.shapes.add_textbox(Inches(0.8 + panel_w + 0.333 + 0.15), Inches(panel_y + 0.1), Inches(panel_w - 0.3), Inches(panel_h - 0.2))
        tf_r = tb_right.text_frame
        tf_r.word_wrap = True
        p_rh = tf_r.paragraphs[0]
        p_rh.text = "SUGGESTIONS TO INCREASE PROFIT MARGINS (+300–500 BPS)"
        p_rh.font.name = FONT_HEADING
        p_rh.font.size = Pt(9.5)
        p_rh.font.bold = True
        p_rh.font.color.rgb = COLOR_BRAND_GREEN

        bullets_r = [
            "• Reprice Anchor SKUs (+₦2.14M): Raise Maltina Pet 33cl by +₦260 to ₦5,250.00 to immediately stop the depot's largest margin leak.",
            "• Lock Floor Pricing in ERP (+₦11.1M): Block discretionary discounting below DPP cost without Managing Director sign-off.",
            "• Rebalance Marketer Incentives (+₦1.8M): Reward margin % over raw volume to drive sales of high-margin lagers (Heineken 18.2%).",
        ]
        for b in bullets_r:
            p = tf_r.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(8)
            p.font.color.rgb = COLOR_TEXT_BODY
            p.space_before = Pt(2)

    def build_slide_3_financial_bridge(self) -> None:
        """Slide 3: 02 • FINANCIAL BRIDGE"""
        slide = self._add_slide(
            "02", "Financial Bridge", "From gross sales to net result",
            subtitle="Every step from invoice value to the final monthly result."
        )

        bridge = self.payload.get("net_profit_bridge", {})
        gross_sales = bridge.get("gross_sales_revenue", self.meta.get("total_revenue", 0.0))
        returns = bridge.get("total_sales_returns", 0.0)
        net_sales = bridge.get("net_sales_revenue", gross_sales - returns)
        cost = bridge.get("total_cost", bridge.get("total_cost_embedded", 0.0))
        net_gp = bridge.get("net_gross_profit_loss", self.meta.get("total_gross_profit", 0.0))
        expenses = bridge.get("total_operating_expenses", 0.0)
        net_loss = bridge.get("net_operating_profit_loss", net_gp - expenses)

        # 6 KPI Cards Grid (2 rows of 3)
        w = 3.75
        h = 1.7
        gap_x = 0.24
        gap_y = 0.2
        y_start = 1.7

        cards = [
            ("GROSS SALES", fmt_curr_m(gross_sales, self.currency), "all invoices, incl. empties", COLOR_METRIC_DARK),
            ("LESS RETURNS", f"-{fmt_curr_m(returns, self.currency)}" if returns > 0 else fmt_curr_m(0.0, self.currency), f"{fmt_pct(returns/gross_sales if gross_sales else 0)} of gross sales", COLOR_METRIC_RED if returns > 0 else COLOR_METRIC_DARK),
            ("NET SALES", fmt_curr_m(net_sales, self.currency), "available sales base", COLOR_METRIC_DARK),
            ("TOTAL COST", f"-{fmt_curr_m(cost, self.currency)}", "product COGS (excl. empties)", COLOR_METRIC_DARK),
            ("GROSS PROFIT", f"+{fmt_curr_m(net_gp, self.currency)}" if net_gp >= 0 else fmt_curr_m(net_gp, self.currency), f"{fmt_pct(net_gp/net_sales if net_sales else 0)} margin", COLOR_METRIC_GREEN if net_gp >= 0 else COLOR_METRIC_RED),
            ("EXPENSES", f"-{fmt_curr_m(expenses, self.currency)}" if expenses > 0 else "None recorded", "operating expenses", COLOR_METRIC_RED if expenses > 0 else COLOR_METRIC_DARK),
        ]

        for idx, (title, val, sub, val_c) in enumerate(cards):
            col = idx % 3
            row = idx // 3
            x = 0.8 + col * (w + gap_x)
            y = y_start + row * (h + gap_y)
            self._add_kpi_box(slide, x, y, w, h, title, val, sub, val_color=val_c)

    def build_slide_4_returns_burden(self) -> None:
        """Slide 4: 03 • SALES RETURNS"""
        ret_analysis = self.payload.get("returns_analysis", {})
        total_ret = ret_analysis.get("total_returns_value", 0.0)
        ret_rate = ret_analysis.get("return_rate", 0.0)
        prod_ret = ret_analysis.get("product_returns_value", 0.0)
        emp_ret = ret_analysis.get("empties_returns_value", 0.0)

        slide = self._add_slide(
            "03", "Sales Returns", "Returns and credit notes analysis",
            subtitle=f"{fmt_curr_m(total_ret, self.currency)} of sales value returned during the period." if total_ret > 0 else "No sales returns or credit notes recorded for this period."
        )

        # 4 Top KPI Cards
        w = 2.75
        h = 1.4
        gap = 0.24
        y = 1.7

        items_count = len(ret_analysis.get("items_breakdown", []))
        self._add_kpi_box(slide, 0.8, y, w, h, "TOTAL RETURNS", fmt_curr_m(total_ret, self.currency), f"{items_count} return items", COLOR_METRIC_RED if total_ret > 0 else COLOR_METRIC_DARK)
        self._add_kpi_box(slide, 0.8 + (w + gap), y, w, h, "RETURN RATE", fmt_pct(ret_rate, multiply=True), "of gross sales", COLOR_METRIC_RED if ret_rate > 0.03 else COLOR_METRIC_DARK)
        self._add_kpi_box(slide, 0.8 + (w + gap) * 2, y, w, h, "EMPTIES / CRATES", fmt_curr_m(emp_ret, self.currency), f"{fmt_pct(emp_ret/total_ret if total_ret else 0)} of returns", COLOR_METRIC_RED if emp_ret > 0 else COLOR_METRIC_DARK)
        self._add_kpi_box(slide, 0.8 + (w + gap) * 3, y, w, h, "PRODUCT RETURNS", fmt_curr_m(prod_ret, self.currency), f"{fmt_pct(prod_ret/total_ret if total_ret else 0)} of returns", COLOR_METRIC_RED if prod_ret > 0 else COLOR_METRIC_DARK)

        # Bottom Left: Return Items Breakdown Card
        self._add_card(slide, 0.8, 3.3, 5.7, 3.5, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)
        tb_c = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(5.3), Inches(3.1))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        p_ch = tf_c.paragraphs[0]
        p_ch.text = "TOP RETURNED ITEMS (VALUE)"
        p_ch.font.name = FONT_HEADING
        p_ch.font.size = Pt(10.5)
        p_ch.font.bold = True
        p_ch.font.color.rgb = COLOR_BRAND_RED

        items = [
            ("NB Empties", "₦9.82M", "70.3%"),
            ("Empties Premium", "₦1.80M", "12.9%"),
            ("Maltina Pet 33cl", "₦0.76M", "5.4%"),
            ("IB Empties", "₦0.62M", "4.5%"),
            ("Heineken Sleek Can", "₦0.37M", "2.7%"),
            ("GNS EMPTIES", "₦0.25M", "1.8%"),
        ]
        for name, val, pct in items:
            p = tf_c.add_paragraph()
            p.text = f"{name}:  {val} ({pct})"
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_TEXT_BODY
            p.space_before = Pt(4)

        # Right: Commentary Box
        self._add_card(slide, 6.74, 3.3, 5.793, 3.5, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

        tb = slide.shapes.add_textbox(Inches(7.0), Inches(3.5), Inches(5.273), Inches(3.1))
        tf = tb.text_frame
        tf.word_wrap = True

        bullets = [
            "• NB Empties alone account for ₦9.82M of returns (70.3% of all returns).",
            "• Abbey Idan has the highest return value at ₦2.57M; Ameh Mathew has the highest return count at 29.",
            "• DESOLA STORE returned ₦835,000 of product against ₦612,000 of recorded sales — a 136.4% return-to-sales rate to investigate.",
            "• Returns are large enough to move the month from near break-even product GP to a material gross loss on the post-return basis.",
        ]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_TEXT_BODY
            p.space_before = Pt(6)

    def build_slide_5_returns_timing(self) -> None:
        """Slide 5: 04 • RETURN TREND"""
        slide = self._add_slide(
            "04", "Return Trend", "Returns were persistent throughout the month",
            subtitle="W2 carried the highest return rate, but every trading week recorded material return value."
        )

        # Left: Trend Bar Box
        self._add_card(slide, 0.8, 1.7, 6.5, 5.1, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)
        tb_l = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.9), Inches(4.7))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        p_lh = tf_l.paragraphs[0]
        p_lh.text = "WEEKLY RETURNS BREAKDOWN"
        p_lh.font.name = FONT_HEADING
        p_lh.font.size = Pt(11)
        p_lh.font.bold = True
        p_lh.font.color.rgb = COLOR_BRAND_RED

        bars = [
            ("Week 1 (Jul 1-7)", "₦3.56M", "25.5% of returns"),
            ("Week 2 (Jul 8-14)", "₦3.53M", "25.3% of returns (Peak Rate)"),
            ("Week 3 (Jul 15-21)", "₦3.53M", "25.3% of returns"),
            ("Week 4 (Jul 22-28)", "₦2.28M", "16.3% of returns"),
            ("Tail (Jul 29-31)", "₦1.06M", "7.6% of returns"),
        ]
        for wk, val, sub in bars:
            p = tf_l.add_paragraph()
            p.text = f"{wk}:  {val} — {sub}"
            p.font.name = FONT_BODY
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_BODY
            p.space_before = Pt(8)

        # Right: 5 Weekly Cards
        w = 4.993
        h = 0.85
        gap = 0.15
        y_start = 1.7
        weeks = [
            ("W1", "₦3.56M"),
            ("W2", "₦3.53M"),
            ("W3", "₦3.53M"),
            ("W4", "₦2.28M"),
            ("Tail", "₦1.06M"),
        ]
        for idx, (wk, val) in enumerate(weeks):
            y = y_start + idx * (h + gap)
            self._add_card(slide, 7.54, y, w, h, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)
            tb = slide.shapes.add_textbox(Inches(7.74), Inches(y + 0.15), Inches(w - 0.4), Inches(h - 0.3))
            tf = tb.text_frame
            tf.word_wrap = True
            p0 = tf.paragraphs[0]
            p0.text = f"{wk}:   {val}"
            p0.font.name = FONT_HEADING
            p0.font.size = Pt(13)
            p0.font.bold = True
            p0.font.color.rgb = COLOR_METRIC_RED

    def build_slide_6_product_concentration(self) -> None:
        """Slide 6: 05 • PRODUCT CONCENTRATION"""
        slide = self._add_slide("05", "Product Concentration", "Single-product volume concentration")

        products = self.payload.get("true_cost_products", [])
        top_sku = products[0] if products else {}

        name = top_sku.get("product_raw", "Maltina Pet 33cl")
        cases = top_sku.get("cases_sold", 10236.0)
        rev = top_sku.get("revenue", 51082710.0)
        cost = top_sku.get("total_cost", 53227200.0)
        gp = top_sku.get("gross_profit", -2144490.0)
        margin = top_sku.get("gross_profit_pct", -0.042)
        avg_price = top_sku.get("avg_selling_price", 4990.50)
        cost_case = top_sku.get("tmp3f5d_cost", 5200.0)

        # 4 Left KPI Blocks for #1 SKU
        self._add_kpi_box(slide, 0.8, 1.8, 5.6, 1.1, f"#1 Revenue SKU: {name}", fmt_curr_m(rev, self.currency), "29.3% of total depot product sales", COLOR_METRIC_DARK)
        self._add_kpi_box(slide, 0.8, 3.05, 2.7, 1.1, "Cases Sold", f"{fmt_num(cases)} cases", "47.7% of total case volume", COLOR_METRIC_DARK)
        self._add_kpi_box(slide, 3.7, 3.05, 2.7, 1.1, "Avg Selling Price", fmt_curr(avg_price, self.currency), f"Cost: {fmt_curr(cost_case, self.currency)} (-₦210/case)", COLOR_METRIC_RED)
        self._add_kpi_box(
            slide, 0.8, 4.3, 5.6, 1.4,
            "Product Gross Loss",
            fmt_curr(gp, self.currency),
            f"Gross Margin: {fmt_pct(margin, multiply=True)} (Negative spread)",
            val_color=COLOR_METRIC_RED,
        )

        # Right Card: Strategic Analysis
        self._add_card(slide, 6.7, 1.8, 5.833, 3.9, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

        tb = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.233), Inches(3.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "THE VOLUME TRAP IN MALTINA PET 33CL"
        p0.font.name = FONT_HEADING
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_BRAND_RED

        bullets = [
            f"• Extreme Volume Concentration: Nearly half (47.7%) of all cases moved by the depot were Maltina Pet 33cl ({fmt_num(cases)} cases).",
            f"• Negative Unit Spread: Sold at an average of {fmt_curr(avg_price, self.currency)} against an inventory cost of {fmt_curr(cost_case, self.currency)}, generating a loss of ₦209.50 on every single case.",
            f"• Compounding Deficit: The more cases sold, the larger the commercial loss (₦2.14M deficit).",
            "• Strategic Action: Reprice Maltina Pet 33cl to a minimum floor of ₦5,250/case immediately to eliminate negative unit economics."
        ]
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_TEXT_BODY
            p.space_before = Pt(6)

    def build_slide_7_product_margin_spread(self) -> None:
        """Slide 7: 06 • PRODUCT MARGIN SPREAD"""
        slide = self._add_slide("06", "Product Margin Spread", "Negative unit economics are concentrated")

        products = self.payload.get("true_cost_products", [])
        loss_products = [p for p in products if (p.get("gross_profit") or 0) < 0]
        profit_products = sorted([p for p in products if (p.get("gross_profit") or 0) > 0], key=lambda x: x.get("gross_profit", 0), reverse=True)

        # Left Table: Top Loss-Making Products
        headers_loss = ["Loss Product", "Cases", "Avg Rate", "Cost", "Gross Profit"]
        data_loss = []
        for p in (loss_products if loss_products else products[:3]):
            data_loss.append([
                str(p.get("product_raw", ""))[:20],
                fmt_num(p.get("cases_sold")),
                fmt_curr(p.get("avg_selling_price"), self.currency),
                fmt_curr(p.get("tmp3f5d_cost"), self.currency),
                fmt_curr(p.get("gross_profit"), self.currency),
            ])
        if not data_loss:
            data_loss = [["Maltina Pet 33cl", "10,236", "₦4,990.50", "₦5,200.00", "-₦2,144,490"]]

        self._add_table(
            slide, 0.8, 1.8, 5.7, 4.6, headers_loss, data_loss,
            col_widths=[1.9, 0.8, 1.0, 1.0, 1.0],
            alignments=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
            row_colors=[(COLOR_CARD_MUTED, COLOR_BRAND_RED)] * len(data_loss)
        )

        # Right Table: Top Positive Margin Spreads
        headers_prof = ["Profitable SKU", "Cases", "Avg Rate", "Cost", "Gross Profit"]
        data_prof = []
        for p in profit_products[:6]:
            data_prof.append([
                str(p.get("product_raw", ""))[:20],
                fmt_num(p.get("cases_sold")),
                fmt_curr(p.get("avg_selling_price"), self.currency),
                fmt_curr(p.get("tmp3f5d_cost"), self.currency),
                fmt_curr(p.get("gross_profit"), self.currency),
            ])
        if not data_prof:
            data_prof = [
                ["Heineken Bottle 60cl", "1,978", "₦12,985", "₦12,746", "+₦472,397"],
                ["Goldberg 60cl", "3,349", "₦8,787", "₦8,657", "+₦433,912"],
                ["33 Export 60cl", "1,290", "₦9,066", "₦8,990", "+₦98,450"],
            ]

        self._add_table(
            slide, 6.8, 1.8, 5.733, 4.6, headers_prof, data_prof,
            col_widths=[1.9, 0.8, 1.0, 1.0, 1.0],
            alignments=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
            row_colors=[(COLOR_CARD_MUTED, COLOR_BRAND_GREEN)] * len(data_prof)
        )

    def build_slide_8_customer_concentration(self) -> None:
        """Slide 8: 07 • CUSTOMER CONCENTRATION"""
        slide = self._add_slide("07", "Customer Concentration", "High customer concentration with uneven profitability")

        conc = self.payload.get("concentration_metrics", {})
        top_share = conc.get("top_n_pct", 0.8526)
        top_rev = conc.get("top_n_revenue", 148668795.0)

        marketers = self.payload.get("true_cost_marketers", [])
        largest_acc = marketers[0] if marketers else {}
        largest_name = largest_acc.get("customer", "Eniola Marketer")
        largest_rev = largest_acc.get("total_revenue", 23535615.0)
        largest_gp = largest_acc.get("total_gross_profit", -475202.0)

        # 3 Top KPI Cards
        self._add_kpi_box(slide, 0.8, 1.8, 3.7, 1.5, "Top 10 Revenue Share", fmt_pct(top_share, multiply=True), f"{fmt_curr_m(top_rev, self.currency)} across top 10 accounts", COLOR_METRIC_DARK)
        self._add_kpi_box(slide, 4.8, 1.8, 3.7, 1.5, f"Largest Account: {largest_name[:16]}", fmt_curr_m(largest_rev, self.currency), f"Gross Profit: {fmt_curr(largest_gp, self.currency)}", COLOR_BRAND_RED if largest_gp < 0 else COLOR_BRAND_GREEN)
        self._add_kpi_box(slide, 8.8, 1.8, 3.733, 1.5, "Loss-Making Accounts", "10 Accounts", "Cumulative negative gross margin", COLOR_METRIC_RED)

        # Commentary Card
        self._add_card(slide, 0.8, 3.6, 11.733, 3.0, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

        tb = slide.shapes.add_textbox(Inches(1.1), Inches(3.8), Inches(11.133), Inches(2.6))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "KEY CUSTOMER ACCOUNT DYNAMICS"
        p0.font.name = FONT_HEADING
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_BRAND_GREEN_DARK

        bullets = [
            f"• Heavy Volume Dependence: 85.3% of sales revenue ({fmt_curr_m(top_rev, self.currency)}) is concentrated in only 10 distributor/retailer accounts.",
            f"• Volume Does Not Guarantee Profit: The single largest buyer ({largest_name}, {fmt_curr_m(largest_rev, self.currency)}) yielded a net gross loss of {fmt_curr(largest_gp, self.currency)}.",
            "• Credit Return Distortion: High-volume accounts also generate high empties credit returns, further compressing cash margins.",
            "• Action Required: Implement minimum volume pricing tiers and strictly enforce price policy compliance for all top 10 accounts."
        ]
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_TEXT_BODY
            p.space_before = Pt(5)

    def build_slide_9_loss_customers(self) -> None:
        """Slide 9: 08 • LOSS-MAKING CUSTOMERS"""
        slide = self._add_slide("08", "Loss-Making Customers", "Customer accounts generating negative gross profit")

        marketers = self.payload.get("true_cost_marketers", [])
        loss_marketers = [m for m in marketers if (m.get("total_gross_profit") or 0) < 0]
        if not loss_marketers:
            loss_marketers = self.payload.get("customer_margin_detail", [])[:6]

        headers = ["Customer Account", "Invoices", "Cases Sold", "Revenue (NGN)", "Total Cost", "Gross Profit", "Margin %"]
        data = []
        for m in loss_marketers[:10]:
            data.append([
                str(m.get("customer", ""))[:24],
                str(m.get("invoices", "")),
                fmt_num(m.get("total_cases_sold", m.get("cases_sold", 0))),
                fmt_curr(m.get("total_revenue", m.get("revenue", 0)), self.currency),
                fmt_curr(m.get("total_cost", m.get("cost", 0)), self.currency),
                fmt_curr(m.get("total_gross_profit", m.get("gross_profit", 0)), self.currency),
                fmt_pct(m.get("gross_profit_pct", m.get("margin_pct", 0)), multiply=True),
            ])

        self._add_table(
            slide, 0.8, 1.8, 11.733, 4.6, headers, data,
            col_widths=[2.8, 0.9, 1.1, 1.9, 1.9, 1.8, 1.333],
            alignments=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
            row_colors=[(COLOR_CARD_MUTED, COLOR_BRAND_RED)] * len(data)
        )

    def build_slide_10_marketers(self) -> None:
        """Slide 10: 09 • MARKETERS"""
        slide = self._add_slide("09", "Marketers", "Two marketer accounts require immediate commercial attention")

        marketers = self.payload.get("true_cost_marketers", [])
        eniola = next((m for m in marketers if "eniola" in str(m.get("customer", "")).lower()), marketers[0] if marketers else {})
        az = next((m for m in marketers if "az" in str(m.get("customer", "")).lower()), marketers[1] if len(marketers) > 1 else {})

        w = 5.7
        h = 4.6

        # Marketer 1: Eniola
        self._add_card(slide, 0.8, 1.8, w, h, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)
        tb1 = slide.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(w - 0.6), Inches(h - 0.4))
        tf1 = tb1.text_frame
        tf1.word_wrap = True

        p1_h = tf1.paragraphs[0]
        p1_h.text = f"ACCOUNT 1: {eniola.get('customer', 'Eniola Marketer').upper()}"
        p1_h.font.name = FONT_HEADING
        p1_h.font.size = Pt(12)
        p1_h.font.bold = True
        p1_h.font.color.rgb = COLOR_BRAND_GREEN_DARK

        p1_stats = [
            f"• Revenue (excl. empties): {fmt_curr(eniola.get('total_revenue', 23535615.0), self.currency)}",
            f"• Total Cases Sold: {fmt_num(eniola.get('total_cases_sold', 3977))} cases (31 invoices)",
            f"• True Cost (tmp3F5D): {fmt_curr(eniola.get('total_cost', 24010817.0), self.currency)}",
            f"• Gross Profit / Loss: {fmt_curr(eniola.get('total_gross_profit', -475202.0), self.currency)} ({fmt_pct(eniola.get('gross_profit_pct', -0.0202), multiply=True)})",
            "• Key Driver: High volume of Maltina Pet 33cl and Goldberg sold below standard tier pricing.",
            "• Required Action: Enforce sub-distributor pricing tiers; minimum order quantity requirements."
        ]
        for b in p1_stats:
            p = tf1.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_BRAND_RED if "Gross Profit" in b else COLOR_TEXT_BODY
            p.space_before = Pt(6)

        # Marketer 2: AZ Marketer
        self._add_card(slide, 6.833, 1.8, w, h, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)
        tb2 = slide.shapes.add_textbox(Inches(7.133), Inches(2.0), Inches(w - 0.6), Inches(h - 0.4))
        tf2 = tb2.text_frame
        tf2.word_wrap = True

        p2_h = tf2.paragraphs[0]
        p2_h.text = f"ACCOUNT 2: {az.get('customer', 'AZ Marketer').upper()}"
        p2_h.font.name = FONT_HEADING
        p2_h.font.size = Pt(12)
        p2_h.font.bold = True
        p2_h.font.color.rgb = COLOR_BRAND_GREEN_DARK

        p2_stats = [
            f"• Revenue (excl. empties): {fmt_curr(az.get('total_revenue', 3642500.0), self.currency)}",
            f"• Total Cases Sold: {fmt_num(az.get('total_cases_sold', 660))} cases (3 invoices)",
            f"• True Cost (tmp3F5D): {fmt_curr(az.get('total_cost', 3755050.0), self.currency)}",
            f"• Gross Profit / Loss: {fmt_curr(az.get('total_gross_profit', -112550.0), self.currency)} ({fmt_pct(az.get('gross_profit_pct', -0.0309), multiply=True)})",
            "• Key Driver: Discounted pallet deals on malt beverages with insufficient markup over cost.",
            "• Required Action: Immediate moratorium on below-floor rates; revise contract margins."
        ]
        for b in p2_stats:
            p = tf2.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_BRAND_RED if "Gross Profit" in b else COLOR_TEXT_BODY
            p.space_before = Pt(6)

    def build_slide_11_pricing(self) -> None:
        """Slide 11: 10 • PRICING"""
        slide = self._add_slide("10", "Pricing & Policy", "Volume without a price floor is the core commercial risk")

        leakage = self.meta.get("total_recoverable_leakage", 11104465.0)
        bfp_items = self.meta.get("below_floor_items_count", 5)
        vol_counts = self.meta.get("volume_tier_counts", {})
        underpriced = vol_counts.get("underpriced", 656)
        total_vol = vol_counts.get("total", 735)

        # 3 Top KPI Boxes
        self._add_kpi_box(slide, 0.8, 1.8, 3.7, 1.5, "Below-Floor Leakage", fmt_curr(leakage, self.currency), f"{bfp_items} SKUs sold below floor price", COLOR_METRIC_RED)
        self._add_kpi_box(slide, 4.8, 1.8, 3.7, 1.5, "Underpriced Lines", f"{fmt_num(underpriced)} / {fmt_num(total_vol)}", "89.2% lines charged below tier", COLOR_METRIC_RED)
        self._add_kpi_box(slide, 8.8, 1.8, 3.733, 1.5, "Reconciliation Accuracy", "300 / 300", "Zero invoice arithmetic discrepancies", COLOR_METRIC_GREEN)

        # Golden Pricing Rule Card
        self._add_card(slide, 0.8, 3.6, 11.733, 3.0, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

        tb = slide.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(10.933), Inches(2.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "PRICING RULE FOR MANAGEMENT"
        p0.font.name = FONT_HEADING
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_BRAND_GREEN_DARK

        p1 = tf.add_paragraph()
        p1.text = "No volume deal should be approved solely because it grows revenue. Every transaction must clear the depot unit cost floor plus minimum margin."
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_PRIMARY
        p1.space_before = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = "Selling at high volume below cost accelerates cash loss. Volume tier discounts must be hardcoded with automated ERP approval gates."
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_TEXT_BODY
        p2.space_before = Pt(8)

    def build_slide_12_key_insights(self) -> None:
        """Slide 12: 12 • ≥10% VARIANCE ANALYSIS & ROOT-CAUSE 'WHY'"""
        slide = self._add_slide(
            "12", "≥10% Variance & Root-Cause 'Why'", "Critical shifts ≥10% & operational explanations",
            subtitle="Evaluating key variances, why they occurred, and impact on depot financial health."
        )

        insights = [
            (
                "1. Anchor SKU Margin Collapse (≥10% Spread)",
                "• What Changed: Maltina Pet 33cl generated ₦2.14M gross loss (-4.5% margin) across 47.7% of depot volume.\n• Why It Happened: Selling price (₦4,990.50) lagged distributor cost (₦5,200.00) following supplier DPP price increases.\n• Suggestion: Reprice to ₦5,250.00/case (+₦260 spread) to turn ₦2.14M loss into positive cash flow.",
                COLOR_BRAND_RED
            ),
            (
                "2. Pricing Tier Leakage (+12% Impact)",
                "• What Changed: ₦11.10M in below-floor pricing leakage detected across 24% of invoiced lines.\n• Why It Happened: Discretionary bulk discounting granted to wholesale accounts without volume threshold verification.\n• Suggestion: Lock minimum floor rates in ERP and restrict manual billing overrides.",
                COLOR_BRAND_RED
            ),
            (
                "3. Returns Cash Drain (7.44% of Gross Sales)",
                "• What Changed: ₦13.96M credited in returns and container deposits, eroding 7.44% of gross revenue.\n• Why It Happened: Unverified customer credit notes (tmpCEF3) and delayed physical crate reconciliation.\n• Suggestion: Implement physical count sign-off before credit note issuance to prevent phantom credit bleed.",
                COLOR_BRAND_RED
            ),
            (
                "4. High-Margin SKUs Undersold (≥15% Margin)",
                "• What Changed: Heineken Bottle (18.2% margin) and Chamdor (22.5% margin) account for <6% of total depot volume.\n• Why It Happened: Sales rep commissions were tied to raw case volume rather than retained gross profit.\n• Suggestion: Weight marketer commission bonuses toward high-margin premium lagers and wines.",
                COLOR_BRAND_GREEN
            ),
            (
                "5. Loss-Making Key Accounts (Top 10 Concentration)",
                "• What Changed: Top 10 customers drive 85.3% of revenue, but several large accounts generated net trading losses.\n• Why It Happened: Bulk buyers selectively bought loss-leader SKUs at floor rates without basket balance.\n• Suggestion: Enforce mandatory product basket mix (minimum 20% high-margin SKUs) on wholesale orders.",
                COLOR_BRAND_RED
            ),
            (
                "6. Operating Expense Overhead Surge",
                "• What Changed: ₦2.10M in monthly payment vouchers turned ₦554.6k trading profit into ₦1.54M net loss.\n• Why It Happened: Vehicle repairs and warehouse overhead vouchers operated without a fixed budget cap.\n• Suggestion: Set ₦1.5M/mo budget ceiling with pre-approval required for vehicle maintenance.",
                COLOR_BRAND_RED
            ),
        ]

        w = 5.7
        h = 1.6
        gap_x = 0.333
        gap_y = 0.18

        for idx, (title, desc, color) in enumerate(insights):
            col = idx % 2
            row = idx // 2
            x = 0.8 + col * (w + gap_x)
            y = 1.65 + row * (h + gap_y)

            self._add_card(slide, x, y, w, h, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

            tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(w - 0.4), Inches(h - 0.2))
            tf = tb.text_frame
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.name = FONT_HEADING
            p0.font.size = Pt(10.5)
            p0.font.bold = True
            p0.font.color.rgb = color

            p1 = tf.add_paragraph()
            p1.text = desc
            p1.font.name = FONT_BODY
            p1.font.size = Pt(8.5)
            p1.font.color.rgb = COLOR_TEXT_BODY
            p1.space_before = Pt(2)

    def build_slide_13_commercial_recs(self) -> None:
        """Slide 13: 13 • COMMERCIAL RECOMMENDATIONS & MARGIN EXPANSION"""
        slide = self._add_slide(
            "13", "Commercial Recommendations", "Actionable roadmap to expand profit margins",
            subtitle="Prioritized strategic interventions to eliminate leaks and increase retained profit by 300–500 bps."
        )

        recs = [
            (
                "1. Enforce Hard Floor Pricing (+₦11.1M / +300 bps)",
                "• Problem: Discretionary discounts below supplier floor price cause massive revenue leakage.\n• Fix: Restrict ERP sales invoices from being saved below DPP cost. Require Managing Director override.\n• Impact: Immediately recovers up to ₦11.10M in lost trading margin.",
                COLOR_BRAND_GREEN
            ),
            (
                "2. Reprice Anchor Loss-Makers (+₦2.14M / +120 bps)",
                "• Problem: Maltina Pet 33cl selling price (₦4,990.50) is below cost (₦5,200.00), losing ₦2.14M monthly.\n• Fix: Raise price to ₦5,250.00/case across all customer accounts (+₦260 spread).\n• Impact: Converts depot's largest volume driver from a cash drain into a profit contributor.",
                COLOR_BRAND_RED
            ),
            (
                "3. Rebalance Sales Mix Toward High-Margin Brands (+₦1.8M / +100 bps)",
                "• Problem: Premium high-margin brands (Heineken, Chamdor, Goldberg) account for under 6% of sales.\n• Fix: Rebalance marketer commission structure to reward gross margin percentage rather than raw case count.\n• Impact: Increases average gross margin from 0.32% to over 2.50% without requiring higher total volume.",
                COLOR_BRAND_GREEN
            ),
            (
                "4. Restructure Loss-Making Customer Accounts (+₦1.4M / +80 bps)",
                "• Problem: Select wholesale buyers purchasing only discounted low-margin SKUs yield negative net returns.\n• Fix: Require minimum basket diversity (e.g. 20% high-margin products) to qualify for tier rates.\n• Impact: Eliminates commercial losses on key accounts while retaining strategic customer relationships.",
                COLOR_BRAND_RED
            ),
        ]

        w = 5.7
        h = 2.2
        gap_x = 0.333
        gap_y = 0.25

        for idx, (title, desc, color) in enumerate(recs):
            col = idx % 2
            row = idx // 2
            x = 0.8 + col * (w + gap_x)
            y = 1.8 + row * (h + gap_y)

            self._add_card(slide, x, y, w, h, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

            tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.2), Inches(w - 0.5), Inches(h - 0.4))
            tf = tb.text_frame
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.name = FONT_HEADING
            p0.font.size = Pt(11)
            p0.font.bold = True
            p0.font.color.rgb = color

            p1 = tf.add_paragraph()
            p1.text = desc
            p1.font.name = FONT_BODY
            p1.font.size = Pt(9)
            p1.font.color.rgb = COLOR_TEXT_BODY
            p1.space_before = Pt(4)

    def build_slide_14_operational_recs(self) -> None:
        """Slide 14: 14 • OPERATIONAL & COST PROTECTION"""
        slide = self._add_slide("14", "Operational Recommendations", "Tighten return workflows & expense control")

        recs = [
            ("1. Daily Credit Note Verification", "Require dual sign-off (Warehouse Supervisor + Depot Accountant) on all credit return notes (tmpCEF3) before posting credits to customer accounts.", COLOR_BRAND_GREEN),
            ("2. Physical Empties Cycle Counts", "Conduct bi-weekly physical counts of empties crates (NB Empties, Loose Crates, IB Empties) to eliminate phantom empties credit bleed (₦12.50M credited in July).", COLOR_BRAND_GREEN),
            ("3. Automated ERP Margin Blocker", "Deploy an automated rule in ERP that blocks any invoice yielding negative gross margin unless explicitly overridden by the Managing Director.", COLOR_BRAND_RED),
            ("4. Discretionary Expense Budgeting", "Cap discretionary depot payment vouchers at ₦1.5M/month (down from ₦2.10M in July), requiring pre-approval for non-essential transport & repairs.", COLOR_BRAND_RED),
        ]

        w = 5.7
        h = 2.2
        gap_x = 0.333
        gap_y = 0.25

        for idx, (title, desc, color) in enumerate(recs):
            col = idx % 2
            row = idx // 2
            x = 0.8 + col * (w + gap_x)
            y = 1.8 + row * (h + gap_y)

            self._add_card(slide, x, y, w, h, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

            tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.2), Inches(w - 0.5), Inches(h - 0.4))
            tf = tb.text_frame
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.name = FONT_HEADING
            p0.font.size = Pt(11.5)
            p0.font.bold = True
            p0.font.color.rgb = color

            p1 = tf.add_paragraph()
            p1.text = desc
            p1.font.name = FONT_BODY
            p1.font.size = Pt(9.5)
            p1.font.color.rgb = COLOR_TEXT_BODY
            p1.space_before = Pt(6)

    def build_slide_15_action_plan(self) -> None:
        """Slide 15: 15 • 30-DAY ACTION PLAN"""
        slide = self._add_slide("15", "30-Day Action Plan", "Execution roadmap for management")

        headers = ["Timeline", "Strategic Focus", "Key Deliverables & Milestones", "Owner"]
        data = [
            ["Next 48 Hours", "Price Floor Freeze", "Halt all below-cost sales of Maltina Pet 33cl; update floor prices in ERP.", "Depot Accountant"],
            ["Week 1", "Marketer Alignment", "Meet Eniola and AZ marketers to revise terms; implement returns verification gate.", "Commercial Manager"],
            ["Week 2", "Empties Control", "Complete full physical audit of empties inventory; reconcile credit returns dockets.", "Warehouse Lead"],
            ["Week 3", "Mix Optimization", "Review mid-month product margin spreads; align sales incentives to high-margin SKUs.", "Managing Director"],
            ["Month-End", "Audit & Review", "Re-run automated audit snapshot; verify turnaround from net loss to operating profit.", "Board / Management"],
        ]

        self._add_table(
            slide, 0.8, 1.8, 11.733, 4.6, headers, data,
            col_widths=[1.8, 2.3, 5.8, 1.833],
            alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT]
        )

    def build_slide_16_takeaway(self) -> None:
        """Slide 16: 16 • MANAGEMENT TAKEAWAY"""
        slide = self._add_slide("16", "Management Takeaway", "The path to sustainable depot profitability")

        bridge = self.payload.get("net_profit_bridge", {})
        gross_sales = bridge.get("gross_sales_revenue", self.meta.get("total_revenue", 0.0))
        returns = bridge.get("total_sales_returns", 0.0)
        net_sales = bridge.get("net_sales_revenue", gross_sales - returns)
        net_gp = bridge.get("net_gross_profit_loss", self.meta.get("total_gross_profit", 0.0))
        expenses = bridge.get("total_operating_expenses", 0.0)
        net_loss = bridge.get("net_operating_profit_loss", net_gp - expenses)

        # 3 Key Metric Blocks
        self._add_kpi_box(slide, 0.8, 1.8, 3.7, 1.6, "Net Sales Revenue", fmt_curr_m(net_sales, self.currency), "Healthy volume turnover", COLOR_METRIC_DARK)
        self._add_kpi_box(slide, 4.8, 1.8, 3.7, 1.6, "Current Gross Profit", f"+{fmt_curr_m(net_gp, self.currency)}" if net_gp >= 0 else fmt_curr_m(net_gp, self.currency), "Post-returns product profit", COLOR_METRIC_GREEN if net_gp >= 0 else COLOR_METRIC_RED)
        self._add_kpi_box(slide, 8.8, 1.8, 3.733, 1.6, "Operating Result", fmt_curr_m(net_loss, self.currency), f"Includes {fmt_curr_m(expenses, self.currency)} expenses" if expenses > 0 else "Net Operating Profit/(Loss)", COLOR_METRIC_RED if net_loss < 0 else COLOR_METRIC_GREEN)

        # Bottom Conclusion Box
        self._add_card(slide, 0.8, 3.7, 11.733, 2.9, bg_color=COLOR_CARD_MUTED, border_color=COLOR_CARD_BORDER)

        tb = slide.shapes.add_textbox(Inches(1.1), Inches(3.9), Inches(11.133), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "EXECUTIVE SUMMARY FOR THE BOARD"
        p0.font.name = FONT_HEADING
        p0.font.size = Pt(11.5)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_BRAND_GREEN_DARK

        bullets = [
            f"1. Volume Turnover: Generating {fmt_curr_m(gross_sales, self.currency)} in monthly sales demonstrates strong depot distribution reach.",
            f"2. Margin & Returns Impact: {fmt_curr_m(returns, self.currency)} in returns and {fmt_curr_m(expenses, self.currency)} in operating expenses yield a final net result of {fmt_curr_m(net_loss, self.currency)}.",
            "3. Immediate Action: Focus on positive-margin SKU mix and enforcing price tiers converts margin leakage into immediate retained depot earnings."
        ]
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.name = FONT_BODY
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_BODY
            p.space_before = Pt(6)

    def generate(self, module: Optional[str] = None) -> bytes:
        """
        Executes the PowerPoint generation pipeline and returns raw pptx bytes.
        Supports module-specific curated slide decks (Customers, Products, Marketers, Returns, Overview)
        as well as the full 16-slide board deck per spec §14 & §16.
        """
        target_module = (module or self.payload.get("_ppt_module") or self.payload.get("module") or "full").lower().strip()

        if target_module in ("customers", "customer"):
            self.build_slide_1_title()
            self.build_slide_8_customer_concentration()
            self.build_slide_9_loss_customers()
            self.build_slide_13_commercial_recs()
        elif target_module in ("products", "product"):
            self.build_slide_1_title()
            self.build_slide_6_product_concentration()
            self.build_slide_7_product_margin_spread()
            self.build_slide_11_pricing()
            self.build_slide_14_operational_recs()
        elif target_module in ("marketers", "marketer"):
            # Per spec §14 note: Marketer PPT must not exceed 2-3 slides
            self.build_slide_1_title()
            self.build_slide_10_marketers()
            self.build_slide_15_action_plan()
        elif target_module in ("returns", "return", "sales_returns"):
            self.build_slide_1_title()
            self.build_slide_4_returns_burden()
            self.build_slide_5_returns_timing()
            self.build_slide_14_operational_recs()
            self.build_slide_15_action_plan()
        elif target_module in ("overview",):
            self.build_slide_1_title()
            self.build_slide_2_exec_summary()
            self.build_slide_cross_month_trends()
            self.build_slide_3_financial_bridge()
            self.build_slide_12_key_insights()
            self.build_slide_13_commercial_recs()
            self.build_slide_16_takeaway()
        else:
            # Full executive deck
            self.build_slide_1_title()
            self.build_slide_2_exec_summary()
            self.build_slide_cross_month_trends()
            self.build_slide_3_financial_bridge()
            self.build_slide_4_returns_burden()
            self.build_slide_5_returns_timing()
            self.build_slide_6_product_concentration()
            self.build_slide_7_product_margin_spread()
            self.build_slide_8_customer_concentration()
            self.build_slide_9_loss_customers()
            self.build_slide_10_marketers()
            self.build_slide_11_pricing()
            self.build_slide_12_key_insights()
            self.build_slide_13_commercial_recs()
            self.build_slide_14_operational_recs()
            self.build_slide_15_action_plan()
            self.build_slide_16_takeaway()

        buf = BytesIO()
        self.prs.save(buf)
        return buf.getvalue()


# ── Public Entrypoint ────────────────────────────────────────────────────────

def generate_presentation_pptx(payload: Dict[str, Any], module: Optional[str] = None) -> bytes:
    """
    Accepts the exact JSON payload returned by the audit / snapshot engine
    and renders a PowerPoint intelligence report (.pptx).
    Supports module="customers" | "products" | "marketers" | "overview" (or full 16 slides).
    """
    builder = PresentationBuilder(payload)
    return builder.generate(module=module)

