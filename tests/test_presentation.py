"""
Unit and functional tests for 16-slide PowerPoint Intelligence Report generator.
"""

from io import BytesIO
import pytest
from pptx import Presentation

from engine.presentation import generate_presentation_pptx, PresentationBuilder


class TestPresentationGenerator:
    """Validates python-pptx presentation generation and slide architecture."""

    def test_presentation_generation_from_snapshot(self, stored_snapshot):
        pptx_bytes = generate_presentation_pptx(stored_snapshot)

        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 40000, "Generated PPTX should be greater than 40KB"

        # Load back with pptx to inspect slide count and structure
        prs = Presentation(BytesIO(pptx_bytes))
        assert len(prs.slides) >= 16, f"Expected at least 16 slides, found {len(prs.slides)}"


        # Validate slide dimensions (16:9 widescreen)
        assert round(prs.slide_width.inches, 2) == 13.33
        assert round(prs.slide_height.inches, 2) == 7.50

    def test_slide_content_and_footers(self, stored_snapshot):
        builder = PresentationBuilder(stored_snapshot)
        pptx_bytes = builder.generate()
        prs = Presentation(BytesIO(pptx_bytes))

        # Check slide 1 (Title slide)
        slide_1 = prs.slides[0]
        text_1 = "\n".join([shape.text_frame.text for shape in slide_1.shapes if shape.has_text_frame])
        assert "JULY 2026 REPORT" in text_1.upper()
        assert "NET SALES" in text_1.upper()

        # Check footer on content slides
        for idx in range(1, 16):
            slide = prs.slides[idx]
            slide_text = "\n".join([shape.text_frame.text for shape in slide.shapes if shape.has_text_frame])
            assert "KANE-JONES" in slide_text
            assert "MANAGEMENT INTELLIGENCE" in slide_text
